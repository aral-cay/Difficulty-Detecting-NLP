import argparse
import os
import pathlib
import pandas as pd
from tqdm import tqdm
import fitz  # PyMuPDF
from pptx import Presentation
from bs4 import BeautifulSoup

def extract_text_pdf(filepath):
    """Extract text from PDF file."""
    try:
        doc = fitz.open(filepath)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting PDF {filepath}: {e}")
        return ""

def extract_text_pptx(filepath):
    """Extract text from PowerPoint file."""
    try:
        prs = Presentation(filepath)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting PPTX {filepath}: {e}")
        return ""

def extract_text_html(filepath):
    """Extract text from HTML file."""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f.read(), 'lxml')
            return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        print(f"Error extracting HTML {filepath}: {e}")
        return ""

def extract_text(filepath):
    """Extract text from file based on extension."""
    ext = pathlib.Path(filepath).suffix.lower()
    if ext == '.pdf':
        return extract_text_pdf(filepath)
    elif ext == '.pptx':
        return extract_text_pptx(filepath)
    elif ext in ['.html', '.htm']:
        return extract_text_html(filepath)
    else:
        return ""

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--downloads", default="data/downloads")
    ap.add_argument("--output", default="data/processed/lecture_texts.tsv")
    args = ap.parse_args()

    os.makedirs(os.path.dirname(args.output), exist_ok=True)
    
    results = []
    
    # Walk through downloads directory (organized by topic_id)
    downloads_path = pathlib.Path(args.downloads)
    if not downloads_path.exists():
        print(f"Error: Downloads directory {args.downloads} does not exist")
        return
    
    files = list(downloads_path.rglob("*"))
    files = [f for f in files if f.is_file() and f.suffix.lower() in ['.pdf', '.pptx', '.html', '.htm']]
    
    for filepath in tqdm(files, desc="Extracting text"):
        # Extract topic_id from parent directory name
        topic_id = filepath.parent.name
        relative_path = str(filepath.relative_to(downloads_path))
        
        text = extract_text(str(filepath))
        if text:
            results.append({
                'topic_id': topic_id,
                'filepath': relative_path,
                'text': text
            })
    
    # Save to TSV
    df = pd.DataFrame(results)
    df.to_csv(args.output, sep='\t', index=False)
    print(f"\nExtracted text from {len(results)} files")
    print(f"Output saved to {args.output}")

if __name__ == "__main__":
    main()

